"""
Generate TAS PPT slides from DB records.
Uses the New_Rev.pptx as a visual template (copies slide XML + image).
"""
import copy
import io
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

TEMPLATES_DIR        = Path(__file__).parent / "templates"
TEMPLATE_PATH        = TEMPLATES_DIR / "AI System 조치 이력 관리_New_Rev.pptx"
LEGACY_TEMPLATE_PATH = TEMPLATES_DIR / "AI System 조치 이력 관리_Legacy_Rev.pptx"
TEMPLATE_SLIDE_IDX   = 5  # first TAS record slide (index 5 = slide 6)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _copy_slide_to_prs(src_slide, dst_prs) -> object:
    """
    Deep-copy a slide (including its image) into a destination Presentation.
    Returns the new slide object.
    """
    blank_layout = dst_prs.slide_layouts[6]  # Blank layout
    new_slide = dst_prs.slides.add_slide(blank_layout)

    sp_tree = new_slide.shapes._spTree
    sp_tree.clear()
    for el in src_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))

    for pic_el in sp_tree.findall(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}pic"
    ):
        parent = pic_el.getparent()
        if parent is not None:
            parent.remove(pic_el)

    for sp_el in list(sp_tree):
        tag = sp_el.tag.split("}")[-1] if "}" in sp_el.tag else sp_el.tag
        if tag == "pic":
            sp_tree.remove(sp_el)

    return new_slide


def _set_cell_text(cell, text: str, preserve_runs: bool = False):
    """Replace all text in a table cell, keeping the first run's formatting."""
    tf = cell.text_frame
    first_run_xml = None
    for para in tf.paragraphs:
        if para.runs:
            first_run_xml = copy.deepcopy(para.runs[0]._r)
            break

    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    lines = text.split("\n") if text else [""]
    for line in lines:
        from pptx.oxml import parse_xml

        p_xml = '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        p_el = parse_xml(p_xml)

        if first_run_xml is not None:
            r_el = copy.deepcopy(first_run_xml)
            t_el = r_el.find(qn("a:t"))
            if t_el is None:
                from lxml import etree
                t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = line
            p_el.append(r_el)
        else:
            from lxml import etree
            r_el = etree.SubElement(p_el, qn("a:r"))
            t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = line

        txBody.append(p_el)


def _fill_table(slide, record: dict):
    """Find main and signature tables in slide and fill with record data."""
    tables = [s for s in slide.shapes if s.shape_type == 19]

    main_tbl = None
    sign_tbl = None
    for t in tables:
        nrows = len(t.table.rows)
        if nrows >= 7:
            main_tbl = t.table
        elif nrows == 4:
            sign_tbl = t.table

    def cell(tbl, r, c):
        return tbl.rows[r].cells[c]

    if main_tbl:
        _set_cell_text(cell(main_tbl, 0, 1), f" Serial No : {record.get('serial_no','')}")
        _set_cell_text(cell(main_tbl, 0, 2), f" Site : {record.get('site','')}")
        _set_cell_text(cell(main_tbl, 0, 3), f" 담당자 : {record.get('manager','')}")

        _set_cell_text(cell(main_tbl, 1, 1), f" □ Issue : {record.get('issue_date','')}")
        _set_cell_text(cell(main_tbl, 1, 2), f" □ 점검 : {record.get('check_date','')}")
        _set_cell_text(cell(main_tbl, 1, 3), f" □ 조치 : {record.get('action_date','')}")

        _set_cell_text(cell(main_tbl, 2, 1), f" □ Core Version : {record.get('core_version','')}")
        _set_cell_text(cell(main_tbl, 2, 2), f" □ Non-Core Version : {record.get('non_core_version','')}")
        _set_cell_text(cell(main_tbl, 2, 3), f" □ H/W 상태 : {record.get('hw_status','')}")

        _set_cell_text(cell(main_tbl, 3, 1), f" {record.get('symptom','')}")
        _set_cell_text(cell(main_tbl, 4, 1), f" {record.get('cause','')}")
        _set_cell_text(cell(main_tbl, 5, 1), f" {record.get('action','')}")
        _set_cell_text(cell(main_tbl, 6, 1), f" {record.get('next_plan','')}")

    if sign_tbl:
        _set_cell_text(cell(sign_tbl, 1, 1), record.get("author_date", ""))
        _set_cell_text(cell(sign_tbl, 2, 1), record.get("author", ""))
        _set_cell_text(cell(sign_tbl, 2, 2), record.get("reviewer", ""))
        _set_cell_text(cell(sign_tbl, 2, 3), record.get("approver", ""))

    for shape in slide.shapes:
        if shape.has_text_frame and "Serial No" in shape.text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = f"Serial No. : {record.get('serial_no','')}"
            break


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def _get_template_slide(system_group: str):
    """Return (src_prs, src_slide) for the given system group."""
    path = LEGACY_TEMPLATE_PATH if system_group == "LEGACY" else TEMPLATE_PATH
    src_prs = Presentation(str(path))
    return src_prs, src_prs.slides[TEMPLATE_SLIDE_IDX]


def generate_single_pptx(record: dict) -> bytes:
    """Generate a single-slide PPTX for one record."""
    grp = record.get("system_group", "NEW")
    src_prs, src_slide = _get_template_slide(grp)

    dst_prs = Presentation()
    dst_prs.slide_width = src_prs.slide_width
    dst_prs.slide_height = src_prs.slide_height

    new_slide = _copy_slide_to_prs(src_slide, dst_prs)
    _fill_table(new_slide, record)

    buf = io.BytesIO()
    dst_prs.save(buf)
    return buf.getvalue()


def generate_multi_pptx(records: list[dict]) -> bytes:
    """
    Generate a multi-slide PPTX for multiple records.
    Records may mix NEW and LEGACY groups; each uses its own template.
    """
    templates: dict[str, tuple] = {}
    for grp in ("NEW", "LEGACY"):
        path = LEGACY_TEMPLATE_PATH if grp == "LEGACY" else TEMPLATE_PATH
        if path.exists():
            prs = Presentation(str(path))
            templates[grp] = (prs, prs.slides[TEMPLATE_SLIDE_IDX])

    if not templates:
        raise RuntimeError("No template files found")

    base_prs_tmpl = templates.get("NEW", list(templates.values())[0])[0]
    dst_prs = Presentation()
    dst_prs.slide_width = base_prs_tmpl.slide_width
    dst_prs.slide_height = base_prs_tmpl.slide_height

    for record in records:
        grp = record.get("system_group", "NEW")
        _, src_slide = templates.get(grp, list(templates.values())[0])
        new_slide = _copy_slide_to_prs(src_slide, dst_prs)
        _fill_table(new_slide, record)

    buf = io.BytesIO()
    dst_prs.save(buf)
    return buf.getvalue()

"""
Parse existing PPTX files into dict records for DB migration.
Supports both Legacy_Rev and New_Rev formats.
"""
import re
from pathlib import Path
from pptx import Presentation


def _strip(text: str) -> str:
    """Remove leading/trailing whitespace and checkbox characters."""
    return text.strip().lstrip("□ ").strip()


def _after_colon(text: str) -> str:
    """Extract value after the first colon."""
    if ":" in text:
        return text.split(":", 1)[1].strip()
    return _strip(text)


def _parse_main_table(table) -> dict:
    """Extract fields from the 7-row main table."""
    data = {}
    rows = table.rows

    def cell(r, c):
        try:
            return rows[r].cells[c].text.strip()
        except IndexError:
            return ""

    data["serial_no"] = _after_colon(cell(0, 1)).replace(" ", "")
    data["site"]      = _after_colon(cell(0, 2))
    data["manager"]   = _after_colon(cell(0, 3))

    data["issue_date"]  = _after_colon(_strip(cell(1, 1)))
    data["check_date"]  = _after_colon(_strip(cell(1, 2)))
    data["action_date"] = _after_colon(_strip(cell(1, 3)))

    core_raw = _strip(cell(2, 1))
    data["core_version"] = _after_colon(core_raw) if "Core Version" in core_raw else core_raw

    noncore_raw = _strip(cell(2, 2))
    data["non_core_version"] = _after_colon(noncore_raw) if "Non-Core" in noncore_raw else noncore_raw

    hw_raw = _strip(cell(2, 3))
    data["hw_status"] = _after_colon(hw_raw) if "H/W" in hw_raw else hw_raw

    data["symptom"]   = cell(3, 1).strip()
    data["cause"]     = cell(4, 1).strip()
    data["action"]    = cell(5, 1).strip()
    data["next_plan"] = cell(6, 1).strip()

    return data


def _parse_sign_table(table) -> dict:
    """Extract author/reviewer/approver from signature table."""
    data = {}
    rows = table.rows

    def cell(r, c):
        try:
            return rows[r].cells[c].text.strip()
        except IndexError:
            return ""

    data["author_date"] = cell(1, 1)
    data["author"]      = cell(2, 1)
    data["reviewer"]    = cell(2, 2)
    data["approver"]    = cell(2, 3)
    return data


def parse_slide(slide) -> dict | None:
    """Parse a single slide. Returns None if not a TAS record slide."""
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame and "Serial No" in shape.text:
            title = shape.text.strip()
            break

    if not title:
        return None

    tables = [s for s in slide.shapes if s.shape_type == 19]
    if not tables:
        return None

    main_table = None
    sign_table = None
    for t in tables:
        nrows = len(t.table.rows)
        if nrows >= 7:
            main_table = t.table
        elif nrows == 4:
            sign_table = t.table

    if main_table is None:
        return None

    data = _parse_main_table(main_table)
    if sign_table:
        data.update(_parse_sign_table(sign_table))

    if not re.match(r"\d{4}-\d+", data.get("serial_no", "")):
        m = re.search(r"(\d{4}-\d+)", title)
        if m:
            data["serial_no"] = m.group(1)
        else:
            return None

    return data


def parse_pptx(pptx_path: str | Path) -> list[dict]:
    """Parse all TAS record slides from a PPTX file. Returns list of dicts."""
    prs = Presentation(str(pptx_path))
    records = []
    seen = set()

    for slide in prs.slides:
        rec = parse_slide(slide)
        if rec and rec["serial_no"] not in seen:
            seen.add(rec["serial_no"])
            records.append(rec)

    return records
